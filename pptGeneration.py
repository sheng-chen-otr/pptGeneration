import os
import sys
import numpy as np
import pandas as pd
from pptx import Presentation
from datetime import date
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from bcParser_v1_0 import *
from pptx.oxml.xmlchemy import OxmlElement
import argparse as argparse 
import glob as glob


parser = argparse.ArgumentParser(prog='PPT Report Generator',description='Generates the powerpoint report for specified trials')

parser.add_argument("-t","--caseNames", metavar='trial',nargs='+',
                    help='Additional trials to be plotted, if doing current trial, leave blank!')
     
args = parser.parse_args()

#### Functions ####
def SubElement(parent, tagname, **kwargs):
    element = OxmlElement(tagname)
    element.attrib.update(kwargs)
    parent.append(element)
    return element


def _set_cell_border(cell, border_color="000000", border_width='500'):
    """ Hack function to enable the setting of border width and border color
        - bottom border only at present
        (c) Steve Canny
    """
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()

    lnR = SubElement(
        tcPr, 'a:lnR', w=border_width, cap='flat', cmpd='sng', algn='ctr')
    solidFill = SubElement(lnR, 'a:solidFill')
    srgbClr = SubElement(solidFill, 'a:srgbClr', val=border_color)
    return cell
################################################################################################

print("#### Running Post Pro Report Generator v2.0 ####")


todays_date = date.today()
path = os.path.split(os.getcwd())[0]
case = os.path.split(os.getcwd())[1]
casePath = os.path.split(path)[0]
job = os.path.split(casePath)[1]




#### TEMPLATE PATH ####
installPath = os.path.dirname(os.path.realpath(__file__))

templatePath = os.path.join(installPath,'pptTemplate','XXXXXX-Template_Presentation.pptx')
print("TEMPLATE PATH: %s" % (templatePath))

#Getting number of trials compared
#numTrials = len(sys.argv)
if args.caseNames is None:
	numTrials = 1
else:	
	numTrials = len(args.caseNames) + 1

caseArray = [case]

#making the list of trials compared
print("Getting data for trials:")
print("\t%s" % (case))
if args.caseNames is not None:
	for i in args.caseNames:
		caseArray.append(i)
		print(i)


def main():
	#initializing the arrays for the trial boundary conditions
	inletMagArray = []
	lastTimeArray = []
	yawArray = []
	wheelRotationArray = []
	simTypeArray = []
	turbModelArray = []
	caseSymArray = []

	for trial in caseArray:
		casePath = os.path.join(path,trial)
		print('\n\tGetting BC data for %s...' % (trial))
		if not os.path.isfile("%s/summary.csv" % (casePath)):
			print("\t\tsummary.csv not found... parsing system files...")
			[inletMag,lastTime,yaw,wheelRotation,simType,turbModel]=bcParser(path,trial)
			inletMagArray.append(inletMag)
			lastTimeArray.append(lastTime)
			yawArray.append(yaw)
			wheelRotationArray.append(wheelRotation)
			simTypeArray.append(simType)
			turbModelArray.append(turbModel)
			if "_half" in trial:
				caseSym = "Half Car"
				caseSymArray.append(caseSym)
			else:
				caseSym = "Full Car"
				caseSymArray.append(caseSym)
		else:
			print("\t\tsummary.csv found... importing data...")
			summaryData = pd.read_csv("%s/summary.csv" % (casePath), index_col=0)
			summaryData = summaryData.transpose()
			inletMagArray.append(summaryData['Velocity'].iloc[0])
			lastTimeArray.append(summaryData['Iterations'].iloc[0])
			yawArray.append(summaryData['Yaw'].iloc[0])
			wheelRotationArray.append(summaryData['Moving Ground'].iloc[0])
			simTypeArray.append(summaryData['Simulation Type'].iloc[0])
			turbModelArray.append(summaryData['Turbulence Model'].iloc[0])
			caseSymArray.append(summaryData['Symmetry'].iloc[0])



	

	#### Setting up Presentation ####

	#getting the presentation template
	prs = Presentation(templatePath)

	print("Making title slide...")
	#setting up the title slide
	titleSlideLayout = prs.slide_layouts[0]
	titleSlide = prs.slides.add_slide(titleSlideLayout)
	title = titleSlide.shapes.title
	subtitle = titleSlide.placeholders[1]
	title.text = "%s - JOB NAME" % (job)
	trialsInReport = ' | '.join(caseArray)
	subtitle.text = "%s" % (trialsInReport)


	#### INFO SLIDE SET UP ####
	print("Making info slide...")
	infoTableLayout = prs.slide_layouts[5]
	infoTableSlide = prs.slides.add_slide(infoTableLayout)
	infoTableTitle = infoTableSlide.shapes.title
	infoTableTitle.text = "Trial Setup and BC" #set the title of the trial set up and boundary conditions slide
	table_placeholder = infoTableSlide.shapes[1]
	nCols = numTrials + 1 #number of columns is number of trials plus 1
	shape = table_placeholder.insert_table(rows=7, cols=nCols)
	table = shape.table

	#making row lables
	rowLabelText = ['Trial','Run Type','Velocity (m/s)','Turb. Model','Wheel Rot.','Yaw Angle (deg)','Sym. Cond.']
	#adding text to the rowLabels
	i = 0
	for text in rowLabelText:
		table.cell(i,0).text = text
		i = i + 1

	#adding data for each trials column
	i = 0 
	while i < numTrials:
		table.cell(0,i+1).text = caseArray[i]
		table.cell(1,i+1).text = simTypeArray[i]
		table.cell(2,i+1).text = str(inletMagArray[i])
		table.cell(3,i+1).text = turbModelArray[i]
		table.cell(4,i+1).text = str(wheelRotationArray[i])
		table.cell(5,i+1).text = str(yawArray[i])
		table.cell(6,i+1).text = caseSymArray[i]
		i = i + 1
		

	for j in range(numTrials + 1):
		for i in range(len(rowLabelText)):
			
			cell = table.cell(i,j)
			cell = _set_cell_border(cell)
			table.cell(i,j).fill.solid()
			table.cell(i,j).text_frame.paragraphs[0].font.size = Pt(14)
			table.cell(i,j).text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
			if i == 0:
				table.cell(i,j).fill.fore_color.rgb = RGBColor(192,192,192)
			else:
				table.cell(i,j).fill.fore_color.rgb = RGBColor(225,225,225)


	#### RESULT TABLE ####
	print("Checking for existance of averaged results...")

	coeffs = np.zeros((6,numTrials))
	n = 0
	for trial in caseArray:
		avgFile = "trial%s_AVG_all_coeff.csv" % (trial)
		if os.path.isfile("%s/%s/%s" % (path,trial,avgFile)):
			print("\t\t\t%s found... importing data..." % (avgFile))
			# Defining the columns to read
			avgData = np.loadtxt("%s/%s/%s" % (path,trial,avgFile), 
						dtype='float', 
						comments='#', 
						delimiter=",",
						skiprows = 1,
						usecols=(1,2,3,4,7,8), 
						unpack=False, ndmin=0)
			
			coeffs[:,n] = avgData
		else: 
			print("		Cannot find average data... skipping...")

		n = n + 1
	



	print("Making results slide...")
	resultsTableLayout = prs.slide_layouts[5]
	resultsTableSlide = prs.slides.add_slide(resultsTableLayout)
	resultsTableTitle = resultsTableSlide.shapes.title
	resultsTableTitle.text = "Results" #set the title of the trial set up and boundary conditions slide
	table_placeholder = resultsTableSlide.shapes[1]
	nCols = numTrials + 1 #number of columns is number of trials plus 1
	shape = table_placeholder.insert_table(rows=8, cols=nCols)
	table = shape.table

	#making row lables
	rowLabelText = ['Trial','CdA','ClA','ClfA','ClrA','0.95 CI. - CdA','0.95 CI. - ClA','% Front']
	#adding text to the rowLabels
	i = 0
	for text in rowLabelText:
		table.cell(i,0).text = text
		i = i + 1

	#adding data for each trials column
	i = 0 
	percentFrontArray = []

	while i < numTrials:

		cl = (coeffs[1,i])
		clf = (coeffs[2,i])
		percentFront = 100*clf/cl
		percentFrontArray.append(percentFront)

		i = i + 1
		

	i = 0
	while i < numTrials:
		#naming first row in column as trial name
		table.cell(0,i+1).text = caseArray[i]

		j = 1
		while j < 7:
			#assigning each value a number
			
			if i == 0:
				table.cell(j,i+1).text = str("%0.3f" % (coeffs[j-1,i]))
			else:
				table.cell(j,i+1).text = str("%0.3f" % (coeffs[j-1,i]))
			j = j + 1
		
		table.cell(7,i+1).text = str("%0.1f"%(percentFrontArray[i]))
		
		i = i + 1


	for j in range(numTrials + 1):
		for i in range(len(rowLabelText)):
			
			cell = table.cell(i,j)
			cell = _set_cell_border(cell)
			table.cell(i,j).fill.solid()
			table.cell(i,j).text_frame.paragraphs[0].font.size = Pt(14)
			table.cell(i,j).text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
			if i == 0:
				table.cell(i,j).fill.fore_color.rgb = RGBColor(192,192,192)
			else:
				table.cell(i,j).fill.fore_color.rgb = RGBColor(225,225,225)

	#### Deltas TABLE ####
	if numTrials > 1:
		percentFrontArray = []
		print("Caltulating deltas results...")
		print("Making deltas slide...")
		resultsTableLayout = prs.slide_layouts[5]
		resultsTableSlide = prs.slides.add_slide(resultsTableLayout)
		resultsTableTitle = resultsTableSlide.shapes.title
		resultsTableTitle.text = "Results Delta to %s" % (caseArray[0]) #set the title of the trial set up and boundary conditions slide, naming it for delta
		table_placeholder = resultsTableSlide.shapes[1]
		nCols = numTrials #number of columns is number of trials plus 1 minus one for deltas, as the first trial is the reference
		shape = table_placeholder.insert_table(rows=6, cols=nCols)
		table = shape.table

		#making row lables
		rowLabelText = ['Trial','deltaCdA','deltaClA','deltaClfA','deltaClrA','delta % Front']
		#rowLabelText = ['Trial','deltaCdA','deltaClA','deltaClfA','deltaClrA','0.95 CI. - deltaCdA','0.95 CI. - deltaClA','delta % Front']
		#adding text to the rowLabels
		i = 0
		for text in rowLabelText:
			table.cell(i,0).text = text
			i = i + 1

		#adding data for each trials column
		i = 1
		percentFrontArray = []
		clRef = coeffs[1,0]
		clfRef = coeffs[2,0]
		while i < numTrials:
			cl = (coeffs[1,i])
			clf = (coeffs[2,i])
			percentFront = (100*clf/cl) - (100*clfRef/clRef)
			percentFrontArray.append(percentFront)
			i = i + 1
		i = 1

		while i < numTrials:
			#naming first row in column as trial name
			table.cell(0,i).text = caseArray[i]

			j = 1
			while j < 5:
				#assigning each value a number
				
				# if j == 5:
				# 	table.cell(j,i+1).text = str("%0.3f" % (np.sqrt(coeffs[4,i]**2+coeffs[4,0]**2)))
				# elif j == 6:
				# 	table.cell(j,i+1).text = str("%0.3f" % (np.sqrt(coeffs[5,i]**2+coeffs[5,0]**2)))
				
				table.cell(j,i).text = str("%0.3f" % (coeffs[j-1,i]-coeffs[j-1,0]))
				j = j + 1
			
			table.cell(5,i).text = str("%0.1f"%(percentFrontArray[i-1]))
			
			i = i + 1


		for j in range(numTrials-1):
			for i in range(len(rowLabelText)):
				
				cell = table.cell(i,j)
				cell = _set_cell_border(cell)
				table.cell(i,j).fill.solid()
				table.cell(i,j).text_frame.paragraphs[0].font.size = Pt(14)
				table.cell(i,j).text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
				if i == 0:
					table.cell(i,j).fill.fore_color.rgb = RGBColor(192,192,192)
				else:
					table.cell(i,j).fill.fore_color.rgb = RGBColor(225,225,225)

	
	#### Confidence Plot SLIDES ####
	print("Checking for Confidence Plot image existance...")
	plotArray = ['forceHistory_Cd','forceHistory_Cl']

	# for trial in caseArray:
	# 	print("		Checking in trial%s" % (trial))
	# 	for plot in plotArray	:
	# 		plotImage = "%s_%s.png" % (trial,plot)
	# 		if os.path.isfile("%s/%s/%s" % (path,trial,plotImage)):
	# 			print("			Found %s..." % (plotImage))
	# 		else:
	# 			print("			Cannot find %s... will skip in report..." % (plotImage))

	# print("Making confidence plot slides...")
	# for plot in plotArray:
	# 	for trial in caseArray:
	# 		confPlotLayout = prs.slide_layouts[6]
	# 		confPlotSlide = prs.slides.add_slide(confPlotLayout)
	# 		confPlotSlideTitle = confPlotSlide.shapes.title
	# 		confPlotSlideTitle.text = "%s - %s" % (plot,trial) #set the title of the geom slides
	# 		confPlot_placeholder = confPlotSlide.shapes[1]
	# 		plotImage = "trial%s_%s.png" % (trial,plot)
	# 		if os.path.isfile("%s/%s/%s" % (path,trial,plotImage)):
	# 			insertConfPlotImage = confPlot_placeholder.insert_picture("%s/%s/%s" % (path,trial,plotImage))
	# 		else:
	# 			pass
	
	
	for plot in plotArray:
		toPlotList = ' '.join(caseArray)
		caseList = '_'.join(caseArray)
		plotImage = "%s_%s.png" % (caseList,plot)
		if os.path.isfile("%s/%s/%s" % (path,caseArray[0],plotImage)):
			print("			Found %s..." % (plotImage))
		else:
			print("			Cannot find %s... will attempt to generate..." % (plotImage))
			#generating the confidence plots if not found
			command = "python3.8 %s/postRun.py --forces -s -t %s" % (installPath,toPlotList)
			os.system(command)

	print("Making confidence plot slides...")
	for plot in plotArray:
		confPlotLayout = prs.slide_layouts[6]
		confPlotSlide = prs.slides.add_slide(confPlotLayout)
		confPlotSlideTitle = confPlotSlide.shapes.title
		confPlotSlideTitle.text = "%s" % (plot) #set the title of the geom slides
		confPlot_placeholder = confPlotSlide.shapes[1]
		plotImage = "%s_%s.png" % (caseList,plot)
		if os.path.isfile("%s/%s/%s" % (path,caseArray[0],plotImage)):
			#insertConfPlotImage = confPlot_placeholder.insert_picture("%s/%s/%s" % (path,caseArray[0],plotImage))
			insertConfPlotImage = confPlotSlide.shapes.add_picture("%s/%s/%s" % (path,caseArray[0],plotImage),left=Inches(2.688), top=Inches(1.7),width=Inches(8.06))
		else:
			pass
	

	#### Development Plot SLIDES ####
	print("Creating development plots...")
	if len(caseArray) < 2:
		caseCommandString = ""
		command = "python3.8 %s/binPlotForces_v3_0.py -s -n -i" % (installPath)
	else:
		caseCommandString = " ".join(caseArray[1:])
		command = "python3.8 %s/binPlotForces_v3_0.py -t %s -s -n -i" % (installPath,caseCommandString)
		
	caseString = "_".join(caseArray)
	os.system(command)
	plotArray = ["cd-development_%s" % (caseString),"cl-development_%s" % (caseString)]


	print("		Checking in trial%s" % (trial))
	for plot in plotArray:
		plotImage = "%s.png" % (plot)
		if os.path.isfile("%s/%s/%s" % (path,case,plotImage)):
			print("			Found %s..." % (plotImage))
		else:
			print("			Cannot find %s... will skip in report..." % (plotImage))

	print("Making development plot slides...")
	for plot in plotArray:

		confPlotLayout = prs.slide_layouts[6]
		confPlotSlide = prs.slides.add_slide(confPlotLayout)
		confPlotSlideTitle = confPlotSlide.shapes.title
		confPlotSlideTitle.text = "%s" % (plot) #set the title of the geom slides
		confPlot_placeholder = confPlotSlide.shapes[1]
		plotImage = "%s.png" % (plot)
		if os.path.isfile("%s/%s/%s" % (path,case,plotImage)):
			insertConfPlotImage = confPlotSlide.shapes.add_picture("%s/%s/%s" % (path,caseArray[0],plotImage),left=Inches(2.688), top=Inches(1.7),width=Inches(8.06))
		else:
			pass
	
	

	

	
	#### GEOMETRY SLIDES ####
	print("Checking for geometry image existance...")
	viewsArray = ['Front','FrontLeft','Left','Bottom','RearLeft','Rear']

	# for trial in caseArray:
	# 	print("		Checking in trial%s" % (trial))
	# 	for view in viewsArray:
	# 		geomImage = "%s_Geom_Surface_%s.png" % (trial,view)
	# 		if os.path.isfile("%s/%s/postProcessing/images/Geom_Surface/%s" % (path,trial,geomImage)):
	# 			print("			Found %s..." % (geomImage))
	# 		else:
	# 			print("			Cannot find %s... will skip in report..." % (geomImage))

	# print("Making geometry slides...")

	# for view in viewsArray:
	# 	for trial in caseArray:
	# 		print('\t\t\tChecking in trial%s' % (trial))
	# 		geomLayout = prs.slide_layouts[4]
	# 		geomSlide = prs.slides.add_slide(geomLayout)
	# 		geomSlideTitle = geomSlide.shapes.title
	# 		geomSlideTitle.text = "Geometry - %s - %s" % (trial,view) #set the title of the geom slides
	# 		geom_placeholder = geomSlide.shapes[1]
	# 		geomImage = "%s_Geom_Surface_%s.png" % (trial,view)
	# 		imagePath = glob.glob("%s/%s/postProcessing/images/Geom_Surface/*_%s.png" % (path,trial,view))
	# 		if len(imagePath) > 0:
	# 			print('\t\t\t\tFound %s' % (imagePath[0].split('/')[-1]))
	# 			insertGeomImage = geomSlide.shapes.add_picture(imagePath[0],left=Inches(2.688), top=Inches(1.7),width=Inches(8.06))
	# 		else:
	# 			pass

	insertImages('Geom_Surface',path,viewsArray,caseArray,prs)
	insertImages('CpMean_Surface',path,viewsArray,caseArray,prs)
	insertImages('CfMean_Surface',path,viewsArray,caseArray,prs)
	insertImages('CpPrime2Mean_Surface',path,viewsArray,caseArray,prs)
	insertImages('Geom_Surface',path,viewsArray,caseArray,prs)
	insertImages('Q_isoSurface',path,viewsArray,caseArray,prs)
	outputReport(prs,caseArray,casePath,todays_date)	
	return

	#### CP SLIDES ####
	print("Checking for CP image existance...")
	viewsArray = ['front','frontLeft','left','bottom','rearLeft','rear']

	for trial in caseArray:
		print("		Checking in trial%s" % (trial))
		for view in viewsArray:
			cpImage = "%s_cP_%s.png" % (trial,view)
			if os.path.isfile("%s/%s/postProcessing/images/%s" % (path,trial,cpImage)):
				print("			Found %s..." % (cpImage))
			else:
				print("			Cannot find %s... will skip in report..." % (cpImage))

	print("Making Cp slides...")
	for view in viewsArray:
		for trial in caseArray:
			cpLayout = prs.slide_layouts[4]
			cpSlide = prs.slides.add_slide(cpLayout)
			cpSlideTitle = cpSlide.shapes.title
			cpSlideTitle.text = "Cp Plot - %s - %s" % (trial,view) #set the title of the geom slides
			cp_placeholder = cpSlide.shapes[1]
			cpImage = "%s_cP_%s.png" % (trial,view)
			if os.path.isfile("%s/%s/postProcessing/images/%s" % (path,trial,cpImage)):
				insertCpImage = cp_placeholder.insert_picture("%s/%s/postProcessing/images/%s" % (path,trial,cpImage))
			else:
				pass

	#### CPX SLIDES ####
	print("Checking for CPX image existance...")
	viewsArray = ['front','frontLeft','left','bottom','rearLeft','rear']

	for trial in caseArray:
		print("		Checking in trial%s" % (trial))
		for view in viewsArray:
			cpXImage = "%s_cPx_%s.png" % (trial,view)
			if os.path.isfile("%s/%s/postProcessing/images/%s" % (path,trial,cpXImage)):
				print("			Found %s..." % (cpXImage))
			else:
				print("			Cannot find %s... will skip in report..." % (cpXImage))

	print("Making CpX slides...")
	for view in viewsArray:
		for trial in caseArray:
			cpXLayout = prs.slide_layouts[4]
			cpXSlide = prs.slides.add_slide(cpXLayout)
			cpXSlideTitle = cpXSlide.shapes.title
			cpXSlideTitle.text = "CpX Plot - %s - %s" % (trial,view) #set the title of the geom slides
			cpX_placeholder = cpXSlide.shapes[1]
			cpXImage = "%s_cPx_%s.png" % (trial,view)
			if os.path.isfile("%s/%s/postProcessing/images/%s" % (path,trial,cpXImage)):
				insertCpXImage = cpX_placeholder.insert_picture("%s/%s/postProcessing/images/%s" % (path,trial,cpXImage))
			else:
				pass

	#### CPZ SLIDES ####
	print("Checking for CPZ image existance...")
	viewsArray = ['front','frontLeft','left','bottom','rearLeft','rear']

	for trial in caseArray:
		print("		Checking in trial%s" % (trial))
		for view in viewsArray:
			cpZImage = "%s_cPz_%s.png" % (trial,view)
			if os.path.isfile("%s/%s/postProcessing/images/%s" % (path,trial,cpZImage)):
				print("			Found %s..." % (cpZImage))
			else:
				print("			Cannot find %s... will skip in report..." % (cpZImage))

	print("Making CpZ slides...")
	for view in viewsArray:
		for trial in caseArray:
			cpZLayout = prs.slide_layouts[4]
			cpZSlide = prs.slides.add_slide(cpZLayout)
			cpZSlideTitle = cpZSlide.shapes.title
			cpZSlideTitle.text = "CpZ Plot - %s - %s" % (trial,view) #set the title of the geom slides
			cpZ_placeholder = cpZSlide.shapes[1]
			cpZImage = "%s_cPz_%s.png" % (trial,view)
			if os.path.isfile("%s/%s/postProcessing/images/%s" % (path,trial,cpZImage)):
				insertCpZImage = cpZ_placeholder.insert_picture("%s/%s/postProcessing/images/%s" % (path,trial,cpZImage))
			else:
				pass



	#### UMeanNear SLIDES ####
	print("Checking for UMeanNear image existance...")
	viewsArray = ['frontLeft','left','bottom','rearLeft','rear']

	for trial in caseArray:
		print("		Checking in trial%s" % (trial))
		for view in viewsArray:
			UMeanNearImage = "%s_UMeanNear_%s.png" % (trial,view)
			if os.path.isfile("%s/%s/postProcessing/images/%s" % (path,trial,UMeanNearImage)):
				print("			Found %s..." % (UMeanNearImage))
			else:
				print("			Cannot find %s... will skip in report..." % (UMeanNearImage))

	print("Making UMeanNear slides...")
	for view in viewsArray:
		for trial in caseArray:
			UMeanNearLayout = prs.slide_layouts[4]
			UMeanNearSlide = prs.slides.add_slide(UMeanNearLayout)
			UMeanNearSlideTitle = UMeanNearSlide.shapes.title
			UMeanNearSlideTitle.text = "UMeanNear Plot - %s - %s" % (trial,view) #set the title of the geom slides
			UMeanNear_placeholder = UMeanNearSlide.shapes[1]
			UMeanNearImage = "%s_UMeanNear_%s.png" % (trial,view)
			if os.path.isfile("%s/%s/postProcessing/images/%s" % (path,trial,UMeanNearImage)):
				insertUMeanNearImage = UMeanNear_placeholder.insert_picture("%s/%s/postProcessing/images/%s" % (path,trial,UMeanNearImage))
			else:
				pass



	#### Ctp Iso SLIDES ####
	print("Checking for CTP Iso image existance...")
	viewsArray = ['frontLeft','left','bottom','rearLeft','rear']

	for trial in caseArray:
		print("		Checking in trial%s" % (trial))
		for view in viewsArray:
			CtpIsoImage = "%s_isoCtp_%s.png" % (trial,view)
			if os.path.isfile("%s/%s/postProcessing/images/%s" % (path,trial,CtpIsoImage)):
				print("			Found %s..." % (CtpIsoImage))
			else:
				print("			Cannot find %s... will skip in report..." % (CtpIsoImage))

	print("Making CTPIso slides...")
	for view in viewsArray:
		for trial in caseArray:
			CtpIsoLayout = prs.slide_layouts[4]
			CtpIsoSlide = prs.slides.add_slide(CtpIsoLayout)
			CtpIsoSlideTitle = CtpIsoSlide.shapes.title
			CtpIsoSlideTitle.text = "Ctp = 0 Iso Plot - %s - %s" % (trial,view) #set the title of the geom slides
			CtpIso_placeholder = CtpIsoSlide.shapes[1]
			CtpIsoImage = "%s_isoCtp_%s.png" % (trial,view)
			if os.path.isfile("%s/%s/postProcessing/images/%s" % (path,trial,CtpIsoImage)):
				insertCtpIsoImage = CtpIso_placeholder.insert_picture("%s/%s/postProcessing/images/%s" % (path,trial,CtpIsoImage))
			else:
				pass




def insertImages(imageType,path,viewsArray,caseArray,prs):
	print("Making geometry slides...")

	for view in viewsArray:
		for trial in caseArray:
			print('\t\t\tChecking in trial%s' % (trial))
			geomLayout = prs.slide_layouts[4]
			geomSlide = prs.slides.add_slide(geomLayout)
			geomSlideTitle = geomSlide.shapes.title
			geomSlideTitle.text = "%s - %s - %s" % (imageType.replace('_',' '),trial,view) #set the title of the geom slides
			geom_placeholder = geomSlide.shapes[1]
			geomImage = "%s_%s_%s.png" % (imageType,trial,view)
			imagePath = glob.glob("%s/%s/postProcessing/images/%s/*_%s.png" % (path,trial,imageType,view))
			if len(imagePath) > 0:
				print('\t\t\t\tFound %s' % (imagePath[0].split('/')[-1]))
				insertGeomImage = geomSlide.shapes.add_picture(imagePath[0],left=Inches(2.688), top=Inches(1.7),width=Inches(8.06))
			else:
				pass
	return prs

def outputReport(prs,caseArray,casePath,todays_date):
	reportName = '_'.join(caseArray)
	print("Saving file to: ../../03_reports/%s_report_%s.pptx" % (reportName,todays_date))
	prs.save("../../03_reports/%s_report_%s.pptx" % (reportName,todays_date))


main()
